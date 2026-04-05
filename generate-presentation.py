#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Color palette - Midnight Executive
colors = {
    'primary': RGBColor(27, 42, 74),        # Navy
    'secondary': RGBColor(212, 175, 55),   # Gold
    'accent': RGBColor(255, 255, 255),     # White
    'darkBg': RGBColor(15, 20, 25),        # Very dark navy
    'lightBg': RGBColor(245, 247, 250),    # Light blue-white
    'text': RGBColor(27, 42, 74),          # Dark navy
    'textLight': RGBColor(255, 255, 255),  # White
    'textMuted': RGBColor(100, 116, 139),  # Slate
    'accentBlue': RGBColor(59, 130, 246),  # Sky blue
}

def add_title_slide(title, subtitle):
    """Add a title slide with navy background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['darkBg']

    # Gold accent line at top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['secondary']
    shape.line.color.rgb = colors['secondary']

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = colors['textLight']
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Footer
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 - Confidential"
    p.font.size = Pt(10)
    p.font.color.rgb = colors['textMuted']
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['lightBg']

    # Top gold bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['secondary']
    shape.line.color.rgb = colors['secondary']

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Content
    yPos = 1.2
    for item in content_items:
        if item['type'] == 'bullet':
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(yPos), Inches(8.4), Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item['text']
            p.font.size = Pt(16)
            p.font.color.rgb = colors['text']
            p.level = 0
            yPos += 0.6
        elif item['type'] == 'box':
            # Add colored box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(yPos),
                Inches(8.4), Inches(item.get('height', 0.8))
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = item.get('bgColor', colors['accentBlue'])
            shape.line.color.rgb = item.get('bgColor', colors['accentBlue'])

            # Text in box
            txBox = slide.shapes.add_textbox(
                Inches(0.9), Inches(yPos + 0.05),
                Inches(8.2), Inches(item.get('height', 0.8) - 0.1)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = item['text']
            p.font.size = Pt(item.get('fontSize', 14))
            p.font.color.rgb = colors['textLight']
            p.alignment = PP_ALIGN.LEFT

            yPos += item.get('height', 0.8) + 0.3

def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors['lightBg']

    # Top bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['secondary']
    shape.line.color.rgb = colors['secondary']

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Center divider
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5), Inches(1.1),
        Inches(0.02), Inches(4.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['secondary']
    shape.line.color.rgb = colors['secondary']

    # Left column title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Left column items
    yPos = 1.8
    for item in left_items:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(yPos), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = colors['text']
        yPos += 0.55

    # Right column title
    txBox = slide.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Right column items
    yPos = 1.8
    for item in right_items:
        txBox = slide.shapes.add_textbox(Inches(5.6), Inches(yPos), Inches(4), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = colors['text']
        yPos += 0.55

# ============================================
# SLIDE 1: Title Slide
# ============================================
add_title_slide("ExpertERP", "La plateforme premium de staffing ERP")

# ============================================
# SLIDE 2: Le Problème
# ============================================
add_content_slide("Le Problème", [
    {'type': 'bullet', 'text': 'ESN et cabinets ERP manquent de consultants qualifiés'},
    {'type': 'bullet', 'text': 'Processus de recrutement lent et fragmenté'},
    {'type': 'bullet', 'text': 'Pas de pool centralisé de ressources vérifiées B2B'},
    {'type': 'bullet', 'text': 'Difficulté à accéder aux talents en dehors de son réseau'},
    {'type': 'bullet', 'text': 'Coûts élevés de sourcing et d\'évaluation des candidats'},
])

# ============================================
# SLIDE 3: Notre Solution
# ============================================
add_content_slide("Notre Solution", [
    {'type': 'box', 'text': 'Plateforme B2B qui connecte entreprises partenaires avec consultants ERP vérifiés', 'bgColor': colors['secondary'], 'height': 0.9},
    {'type': 'bullet', 'text': 'Pool de freelances ERP disponibles immédiatement'},
    {'type': 'bullet', 'text': 'Accès aux ressources des firmes partenaires (B2B privé)'},
    {'type': 'bullet', 'text': 'Profils vérifiés et filtrés par module ERP & certifications'},
    {'type': 'bullet', 'text': 'Mise en relation directe et rapide via messagerie sécurisée'},
])

# ============================================
# SLIDE 4: Comment Ça Marche
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['lightBg']

# Top bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(0.12)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Comment Ça Marche"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = colors['primary']

steps = [
    {'num': '1', 'label': 'Inscription', 'desc': 'Créez un compte entreprise en 2 min'},
    {'num': '2', 'label': 'Recherche', 'desc': 'Filtrez par compétences, modules, TJM'},
    {'num': '3', 'label': 'Contact', 'desc': 'Messagerie directe & mise en relation'},
]

xPos = 0.8
for step in steps:
    # Circle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(xPos + 0.8), Inches(1.5),
        Inches(0.6), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['secondary']
    shape.line.color.rgb = colors['secondary']

    # Number in circle
    txBox = slide.shapes.add_textbox(Inches(xPos + 0.8), Inches(1.5), Inches(0.6), Inches(0.6))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = step['num']
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER

    # Label
    txBox = slide.shapes.add_textbox(Inches(xPos + 0.5), Inches(2.3), Inches(1.2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = step['label']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER

    # Description
    txBox = slide.shapes.add_textbox(Inches(xPos + 0.3), Inches(2.8), Inches(1.6), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step['desc']
    p.font.size = Pt(11)
    p.font.color.rgb = colors['textMuted']
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    if step['num'] != '3':
        arrow = slide.shapes.add_connector(1, Inches(xPos + 1.8), Inches(1.8), Inches(xPos + 2.6), Inches(1.8))
        arrow.line.color.rgb = colors['secondary']
        arrow.line.width = Pt(2)

    xPos += 3

# ============================================
# SLIDE 5: Pour les Entreprises Partenaires
# ============================================
add_two_column_slide(
    "Pour les Entreprises Partenaires",
    "Avantages",
    [
        "Accès au pool B2B complet",
        "Import en lot de ressources (Excel)",
        "Visibilité contrôlée de vos consultants",
        "Gestion d'équipe & profils",
        "Support dédié",
    ],
    "Quoi de plus?",
    [
        "Pas de frais d'inscription",
        "Transparence tarifaire (TJM public)",
        "Vérification des profiles",
        "Historique des contacts",
        "Dashboard analytique",
    ]
)

# ============================================
# SLIDE 6: Pool de Compétences
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['lightBg']

# Top bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(0.12)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Pool de Compétences"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = colors['primary']

# ERP Systems Grid
erp_systems = [
    "SAP S/4HANA",
    "Oracle Cloud",
    "Dynamics 365",
    "Workday",
    "Salesforce",
    "NetSuite",
    "Odoo",
    "Infor",
]

col_x = 0.8
row_y = 1.3
col = 0

for system in erp_systems:
    # Box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(col_x), Inches(row_y),
        Inches(2), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['accentBlue']
    shape.line.color.rgb = colors['accentBlue']

    # Text
    txBox = slide.shapes.add_textbox(Inches(col_x), Inches(row_y), Inches(2), Inches(0.5))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = system
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = colors['textLight']
    p.alignment = PP_ALIGN.CENTER

    col += 1
    col_x += 2.2
    if col == 4:
        col = 0
        col_x = 0.8
        row_y += 0.8

# Modules
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "+ Modules : FI/CO, MM/SD, ABAP, HCM, F&O, Finance, RH..."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = colors['textMuted']

# Certifications
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "+ Certifications : SAP S/4HANA, Oracle Cloud, Workday..."
p.font.size = Pt(13)
p.font.italic = True
p.font.color.rgb = colors['textMuted']

# ============================================
# SLIDE 7: Fonctionnalités Clés
# ============================================
add_content_slide("Fonctionnalités Clés", [
    {'type': 'bullet', 'text': 'Dashboard entreprise avec vue 360 des ressources'},
    {'type': 'bullet', 'text': 'Filtres avancés : modules, compétences, TJM, lieu'},
    {'type': 'bullet', 'text': 'Profils détaillés et vérifiés de chaque consultant'},
    {'type': 'bullet', 'text': 'Messagerie sécurisée et directe (plan payant)'},
    {'type': 'bullet', 'text': 'Import Excel en lot de vos ressources'},
    {'type': 'bullet', 'text': 'Visibilité contrôlée : public ou private B2B'},
])

# ============================================
# SLIDE 8: Sécurité & Conformité
# ============================================
add_content_slide("Sécurité & Conformité", [
    {'type': 'bullet', 'text': 'Sessions signées avec token d\'accès sécurisé'},
    {'type': 'bullet', 'text': 'Row Level Security (RLS) Supabase par utilisateur'},
    {'type': 'bullet', 'text': 'Données chiffrées en transit (HTTPS/TLS)'},
    {'type': 'bullet', 'text': 'Respect RGPD : mentions légales, consentement cookies'},
    {'type': 'bullet', 'text': 'Audit de sécurité complet réalisé (avril 2026)'},
    {'type': 'bullet', 'text': 'Authentification multi-niveaux (email, mot de passe)'},
])

# ============================================
# SLIDE 9: Modèle Tarifaire
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['lightBg']

# Top bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(0.12)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Modèle Tarifaire Simple & Transparent"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = colors['primary']

# Free tier box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.2),
    Inches(4.2), Inches(3.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['accent']
shape.line.color.rgb = colors['secondary']
shape.line.width = Pt(1)

# Free title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4.2), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Inscription Gratuite"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = colors['primary']
p.alignment = PP_ALIGN.CENTER

# Free content
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.9), Inches(3.8), Inches(1.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "• Accès au pool freelances\n• Voir les profils\n• Filtres avancés"
p.font.size = Pt(12)
p.font.color.rgb = colors['text']

# Premium tier box
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.2), Inches(1.2),
    Inches(4.2), Inches(3.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Premium title
txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Plan Messagerie"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = colors['primary']
p.alignment = PP_ALIGN.CENTER

# Premium content
txBox = slide.shapes.add_textbox(Inches(5.4), Inches(1.9), Inches(3.8), Inches(1.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "+ Accès pool B2B\n+ Messagerie directe\n+ Billing Stripe intégré"
p.font.size = Pt(12)
p.font.color.rgb = colors['primary']

# Tagline
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Gratuit jusqu'à 3 messages/mois. Abonnement mensuel après."
p.font.size = Pt(12)
p.font.color.rgb = colors['textMuted']
p.alignment = PP_ALIGN.CENTER
p.font.italic = True

# ============================================
# SLIDE 10: État Actuel & Traction
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['lightBg']

# Top bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(0.12)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "État Actuel & Traction"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = colors['primary']

# Stats
stats = [
    {'num': '2', 'label': 'Entreprises pilotes'},
    {'num': '10+', 'label': 'Ressources B2B'},
    {'num': '12+', 'label': 'Freelances vérifiés'},
    {'num': '100%', 'label': 'Uptime (Vercel)'},
]

stat_x = 0.8
for stat in stats:
    # Number
    txBox = slide.shapes.add_textbox(Inches(stat_x), Inches(1.5), Inches(2), Inches(0.7))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = stat['num']
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Label
    txBox = slide.shapes.add_textbox(Inches(stat_x), Inches(2.3), Inches(2), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = stat['label']
    p.font.size = Pt(11)
    p.font.color.rgb = colors['textMuted']
    p.alignment = PP_ALIGN.CENTER

    stat_x += 2.2

# Infrastructure section
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(8.4), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Infrastructure"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = colors['primary']

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Frontend: Vercel | Backend: Supabase PostgreSQL | Auth: Sessions sécurisées"
p.font.size = Pt(12)
p.font.color.rgb = colors['text']

# ============================================
# SLIDE 11: Roadmap
# ============================================
add_content_slide("Roadmap 2026", [
    {'type': 'box', 'text': 'Q2 2026 : Intégration Stripe complète, matching IA, analytics avancés', 'bgColor': colors['accentBlue']},
    {'type': 'box', 'text': 'Q3 2026 : Mobile app responsive, notifications push', 'bgColor': colors['accentBlue']},
    {'type': 'box', 'text': 'Q4 2026 : Partenariats stratégiques, expansion B2B Europe', 'bgColor': colors['accentBlue']},
])

# ============================================
# SLIDE 12: Pourquoi Nous Rejoindre Maintenant
# ============================================
add_content_slide("Pourquoi Nous Rejoindre Maintenant", [
    {'type': 'bullet', 'text': 'Avantage early adopter : façonnez la plateforme avec nous'},
    {'type': 'bullet', 'text': 'Accès gratuit initial avant monétisation complète'},
    {'type': 'bullet', 'text': 'Plateforme live & stable, infrastructure enterprise-grade'},
    {'type': 'bullet', 'text': 'Support dédié pour intégration et déploiement'},
    {'type': 'bullet', 'text': 'Potentiel d\'évaluation avantageuse lors des futurs tours'},
])

# ============================================
# SLIDE 13: Call to Action / Contact
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['darkBg']

# Top gold bar
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(10), Inches(0.08)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

# Main heading
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
tf = txBox.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Prêt à Découvrir ExpertERP?"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = colors['textLight']
p.alignment = PP_ALIGN.CENTER

# CTA Button
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3), Inches(2.8),
    Inches(4), Inches(0.7)
)
shape.fill.solid()
shape.fill.fore_color.rgb = colors['secondary']
shape.line.color.rgb = colors['secondary']

txBox = slide.shapes.add_textbox(Inches(3), Inches(2.8), Inches(4), Inches(0.7))
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = "Planifiez une Démo"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = colors['primary']
p.alignment = PP_ALIGN.CENTER

# Contact info
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "contact@experterpbub.com\n+1 514 XXX-XXXX"
p.font.size = Pt(14)
p.font.color.rgb = colors['accentBlue']
p.alignment = PP_ALIGN.CENTER

# Footer
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ExpertERP - La plateforme de staffing ERP pour les professionnels"
p.font.size = Pt(11)
p.font.color.rgb = colors['textMuted']
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "/sessions/loving-happy-hopper/mnt/Projects--ExpertERPHUB/ExpertERP_Presentation_Entreprises_2026.pptx"
prs.save(output_path)
print(f"Presentation created successfully at: {output_path}")
